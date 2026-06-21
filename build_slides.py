"""
Generateur de la presentation de soutenance (.pptx) - projet IDS Intelligent / ORNL Driver ID.
Retrace le parcours et les problemes rencontres. Adapte du generateur de l'ancien projet.
Sortie : deliverables/Presentation_IDS_Intelligent.pptx     Lancer : python build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(ROOT, "docs", "assets")
OUT = os.path.join(ROOT, "deliverables", "Presentation_IDS_Intelligent.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

NAVY = RGBColor(0x0F, 0x1E, 0x38); ICE = RGBColor(0xCA, 0xDC, 0xFC)
RED = RGBColor(0xC0, 0x39, 0x2B); TEAL = RGBColor(0x16, 0xA0, 0x85)
AMBER = RGBColor(0xE6, 0x7E, 0x22); INK = RGBColor(0x1F, 0x2D, 0x3D)
MUTED = RGBColor(0x6B, 0x7B, 0x8C); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF2, 0xF6, 0xFA); SERIF, SANS = "Cambria", "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = NAVY if dark else WHITE
    return s


def tb(s, l, t, w, h, lines, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
       font=SANS, anchor=MSO_ANCHOR.TOP, space=4):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        txt, opt = ln if isinstance(ln, tuple) else (ln, {})
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opt.get("size", size)); r.font.bold = opt.get("bold", bold)
        r.font.name = opt.get("font", font); r.font.color.rgb = opt.get("color", color)
    return box


def title(s, text, dark=False, sub=None):
    tb(s, 0.7, 0.45, 12.2, 0.9, text, size=32, bold=True, font=SERIF, color=WHITE if dark else NAVY)
    if sub:
        tb(s, 0.72, 1.25, 12.2, 0.5, sub, size=14, color=ICE if dark else MUTED)


def pic(s, name, l, t, w):
    path = os.path.join(ASSETS, name)
    if not os.path.exists(path): return
    iw, ih = Image.open(path).size
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(w * ih / iw))


def card(s, l, t, w, h, fill=TINT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def circle(s, l, t, d, color, glyph, gcolor=WHITE, gsize=18):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); sh.shadow.inherit = False
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph; r.font.size = Pt(gsize); r.font.bold = True
    r.font.color.rgb = gcolor; r.font.name = SANS
    return sh


def stat(s, l, t, w, number, label, color=RED):
    tb(s, l, t, w, 0.9, number, size=40, bold=True, font=SERIF, color=color, align=PP_ALIGN.CENTER)
    tb(s, l, t + 0.8, w, 0.6, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)


def tag(s, l, t, text, color, glyph="!"):
    circle(s, l, t, 0.42, color, glyph, gsize=15)
    tb(s, l + 0.55, t + 0.04, 7.5, 0.4, text, size=14, bold=True, color=color)


# ============================================================================ 1. TITRE
s = slide(dark=True)
circle(s, 0.7, 0.6, 0.7, RED, "!", gsize=24)
tb(s, 0.7, 2.3, 12, 1.8, [
    ("Detecteur d'intrusion intelligent", {"size": 40, "bold": True, "font": SERIF, "color": WHITE}),
    ("pour bus CAN de vehicule lourd", {"size": 40, "bold": True, "font": SERIF, "color": WHITE})])
tb(s, 0.72, 4.2, 11.5, 0.6,
   "Detection par machine learning d'une cyberattaque (CAN / J1939) - dataset ORNL Driver Identification",
   size=16, color=ICE)
tb(s, 0.72, 6.2, 11, 0.8, [
    ("Projet integrateur 2026   -   Encadrant : Rida Khatoun", {"size": 14, "color": WHITE}),
    ("Kenworth T270  -  50 conducteurs  -  attaque rare (1,46 %)", {"size": 12, "color": ICE})])

# ============================================================================ 2. FIL DU PROJET
s = slide(); title(s, "Le fil du projet", sub="Un parcours d'AI Engineer : de la theorie a l'auto-critique")
phases = [("1", "Comprendre", "contexte, donnees, confondeurs", NAVY),
          ("2", "Eviter les pieges", "fuite par conducteur, GPS, temps", AMBER),
          ("3", "Modeliser", "3 chemins, evaluation honnete", TEAL),
          ("4", "Decouvrir", "la signature d'injection (silence CAN0)", RED),
          ("5", "Eprouver", "attaquant adaptatif, limites, demo", NAVY)]
y = 2.1
for num, head, txt, col in phases:
    circle(s, 0.9, y, 0.55, col, num, gsize=20)
    tb(s, 1.7, y - 0.02, 4.3, 0.5, head, size=18, bold=True, color=INK)
    tb(s, 6.2, y + 0.05, 6.4, 0.5, txt, size=14, color=MUTED)
    y += 0.95
tb(s, 0.9, 6.9, 11.5, 0.5,
   "Regle d'or : chaque brique = un concept explique + une experience mesuree + un verdict assume.",
   size=13, color=TEAL, font=SERIF)

# ============================================================================ 3. ATTAQUE
s = slide(); title(s, "Une cyberattaque sur un camion",
                   sub="Le bus CAN : fiable en temps reel, mais sans aucune securite")
steps = [("1", "Drive-by Wi-Fi", "Un vehicule proche se connecte a l'ELD (boitier reglementaire) par Wi-Fi.", RED),
         ("2", "Firmware remplace", "Reecrit a distance, l'ELD peut desormais ecrire sur le bus CAN.", AMBER),
         ("3", "Injection CAN", "Il injecte une attaque sur l'afficheur : tachymetre/compteur mis a zero.", TEAL)]
y = 2.0
for num, head, txt, col in steps:
    circle(s, 0.8, y, 0.55, col, num, gsize=20)
    tb(s, 1.6, y - 0.05, 6.2, 0.5, head, size=18, bold=True, color=INK)
    tb(s, 1.6, y + 0.45, 6.2, 0.9, txt, size=13, color=MUTED)
    y += 1.45
card(s, 8.4, 2.0, 4.2, 4.3)
tb(s, 8.7, 2.3, 3.6, 0.5, "Le dispositif cle", size=16, bold=True, color=NAVY)
tb(s, 8.7, 2.9, 3.6, 3.2, [
    "50 conducteurs, 3 groupes d'AWARENESS (non prevenu / prevenu / prevenu+se garer).",
    "", "L'attaque survient TOUJOURS au meme lieu et au meme moment.",
    "", "-> ce dispositif est la source des confondeurs."], size=13, color=INK, space=6)

# ============================================================================ 4. PROBLEMATIQUE
s = slide(); title(s, "Problematique et objectifs",
                   sub="La vraie difficulte n'est pas le modele, ce sont les confondeurs")
tb(s, 0.8, 1.9, 5.7, 0.5, "La question piegee", size=18, bold=True, color=RED)
tb(s, 0.8, 2.4, 5.7, 1.4, "L'attaque etant toujours au meme lieu/moment, un modele naif detecte "
   "l'ENDROIT ou l'INSTANT, pas l'attaque. Tout l'enjeu est de ne pas se faire pieger.", size=15, color=INK)
tb(s, 0.8, 4.0, 5.7, 0.5, "Deux angles", size=15, bold=True, color=NAVY)
tb(s, 0.8, 4.5, 5.7, 1.2, ["Classification supervisee : normal ou attaque ?",
   "Detection d'anomalie : ecart au normal (sans labels)"], size=13, color=INK, space=8)
card(s, 7.0, 1.9, 5.6, 4.6)
tb(s, 7.3, 2.2, 5.0, 0.5, "Objectifs (consigne)", size=16, bold=True, color=NAVY)
tb(s, 7.3, 2.8, 5.0, 3.6, ["Explorer et caracteriser normal vs suspect",
   "Pretraiter et selectionner les features", "Comparer plusieurs algorithmes",
   "Decouper train / test (par conducteur)", "Entrainer sur normal et malveillant",
   "Evaluer (precision, rappel, ROC) et optimiser"], size=14, color=INK, space=8)

# ============================================================================ 5. DATASET
s = slide(); title(s, "Le jeu de donnees", sub="ORNL Driver Identification - une flotte sous attaque")
stat(s, 0.8, 2.3, 2.8, "155 902", "fenetres de 1 s", RED)
stat(s, 3.8, 2.3, 2.8, "50", "conducteurs", NAVY)
stat(s, 0.8, 4.1, 2.8, "1,46 %", "taux d'attaque", AMBER)
stat(s, 3.8, 4.1, 2.8, "337", "signaux CAN", NAVY)
card(s, 7.2, 2.3, 5.4, 3.5)
tb(s, 7.5, 2.6, 4.8, 0.5, "Trois modalites", size=16, bold=True, color=NAVY)
tb(s, 7.5, 3.2, 4.8, 2.4, ["CAN (J1939) : moteur, vehicule -> features de travail",
   "Biometrie (HR/EDA, Empatica) : la reaction physiologique",
   "GPS/inertie (VBOX) : CONFONDEUR de lieu -> exclu"], size=14, color=INK, space=8)

# ============================================================================ 6-8. EDA
s = slide(); title(s, "Exploration (1/3) : le confondeur de LIEU",
                   sub="L'attaque est concentree geographiquement -> le GPS triche")
pic(s, "eda_confounder_gps.png", 0.8, 1.9, 7.3)
card(s, 8.5, 1.9, 4.1, 4.2)
tb(s, 8.8, 2.2, 3.5, 3.6, [("d = 2,63", {"size": 34, "bold": True, "color": RED, "font": SERIF}),
   ("L'attaque est 7 a 44x plus concentree que le trajet.", {"size": 13, "color": INK}),
   ("", {}), ("Un modele GPS detecterait l'ENDROIT, pas l'attaque. -> on EXCLUT les 312 features "
   "GPS/inertie.", {"size": 13, "color": INK})], space=6)

s = slide(); title(s, "Exploration (2/3) : le confondeur de TEMPS",
                   sub="Les signaux a derive lente correlent avec le moment de l'attaque")
pic(s, "eda_cohens_d.png", 0.8, 2.0, 7.2)
card(s, 8.4, 2.0, 4.2, 4.1)
tb(s, 8.7, 2.3, 3.6, 3.4, ["L'attaque survient a ~60 % du trajet.", "",
   ("Test CAN_STABLE", {"size": 15, "bold": True, "color": TEAL}),
   ("Exclure les 21 signaux time-drift : 0,632 -> 0,630 (inchange).", {"size": 13, "color": INK}),
   ("Le signal CAN est GENUINE, pas un confondeur temps.", {"size": 13, "color": INK})], space=6)

s = slide(); title(s, "Exploration (3/3) : la reponse biometrique",
                   sub="Le coeur reagit a l'attaque - mais faiblement")
pic(s, "eda_biometric_hr.png", 0.8, 2.0, 7.0)
card(s, 8.2, 2.0, 4.4, 4.1)
tb(s, 8.5, 2.3, 3.8, 3.4, ["La frequence cardiaque monte (+1,7 a +3,1 bpm), max pour le groupe "
   "averti.", "", ("Effet REEL mais modeste (d = 0,14)", {"size": 14, "bold": True, "color": AMBER}),
   ("", {}), ("Une piste d'analyse, pas un detecteur (confirme en Vague 2).", {"size": 13, "color": INK})], size=13, color=INK, space=6)

# ============================================================================ 9. LE PIEGE CENTRAL
s = slide(); title(s, "Le piege central : la fuite par conducteur",
                   sub="Quand un score brillant ne prouve rien")
pic(s, "p2_confounders_demo.png", 0.8, 2.0, 7.0)
card(s, 8.2, 2.0, 4.4, 4.3, fill=NAVY)
tb(s, 8.5, 2.3, 3.8, 3.8, [("0,985 -> 0,632", {"size": 24, "bold": True, "color": RED, "font": SERIF}),
   ("Split aleatoire : 0,985 (le modele reconnait le CONDUCTEUR, pas l'attaque).", {"size": 13, "color": WHITE}),
   ("", {}), ("Split par conducteur : 0,632 (la vraie difficulte).", {"size": 13, "color": ICE}),
   ("", {}), ("Avec GPS : 0,835 mais c'est le LIEU.", {"size": 13, "color": WHITE})], space=6)
tb(s, 0.8, 6.3, 11.6, 0.6, "Garde-fous pour tout le projet : split PAR CONDUCTEUR + CAN seul + "
   "metrique PR-AUC.", size=13, color=TEAL, font=SERIF)

# ============================================================================ 10. MODALITES
s = slide(); title(s, "Features : trois modalites, un seul choix honnete",
                   sub="Ablation (Gradient Boosting, par conducteur)")
rows = [("CAN (337)", "0,756", "travail", TEAL), ("CAN + bio (344)", "0,749", "n'aide pas", MUTED),
        ("Biometrie (7)", "0,014", "= hasard", AMBER), ("GPS (312)", "0,835", "CONFONDEUR", RED)]
x = 0.8
for head, val, role, col in rows:
    card(s, x, 2.1, 2.85, 2.2)
    tb(s, x + 0.2, 2.35, 2.5, 0.5, head, size=14, bold=True, color=INK)
    tb(s, x + 0.2, 2.95, 2.5, 0.7, val, size=26, bold=True, color=col, font=SERIF)
    tb(s, x + 0.2, 3.7, 2.5, 0.4, role, size=12, color=col)
    x += 3.0
tb(s, 0.8, 4.7, 11.8, 1.2, "Le GPS 'bat' le CAN (0,835 > 0,756) mais c'est un piege (le lieu). La "
   "biometrie seule est au niveau du hasard. -> on travaille sur le CAN seul.", size=15, color=INK)

# ============================================================================ 11. CHEMINS
s = slide(); title(s, "Trois chemins de modelisation",
                   sub="Explorer plusieurs pistes plutot qu'un seul modele")
chemins = [("A", "Supervise", "LogReg, SVM, RF, Gradient Boosting", RED),
           ("B", "Anomalie", "Isolation Forest, One-Class SVM, PCA", TEAL),
           ("C", "Deep", "MLP tabulaire, GRU temporel (PyTorch)", AMBER)]
pos = [(0.8, 2.2), (0.8, 4.4), (6.9, 2.2)]
for (lbl, head, models, col), (lx, ly) in zip(chemins, pos):
    card(s, lx, ly, 5.6, 2.0)
    circle(s, lx + 0.25, ly + 0.3, 0.6, col, lbl, gsize=22)
    tb(s, lx + 1.05, ly + 0.35, 4.3, 0.6, head, size=17, bold=True, color=INK)
    tb(s, lx + 0.3, ly + 1.2, 5.0, 0.7, models, size=12, color=MUTED)
card(s, 6.9, 4.4, 5.6, 2.0, fill=NAVY)
tb(s, 7.2, 4.7, 5.0, 1.5, [("Verdict", {"size": 16, "bold": True, "color": ICE}),
   ("Arbres boostes = champion (0,756). Deep 0,54-0,57. Anomalie ~0,02.", {"size": 14, "color": WHITE})], space=6)

# ============================================================================ 12. CHEMIN A
s = slide(); title(s, "Chemin A : les arbres boostes dominent",
                   sub="Comparaison supervisee, validation par conducteur")
pic(s, "p4a_supervised.png", 0.8, 2.0, 7.2)
card(s, 8.4, 2.0, 4.2, 4.3)
tb(s, 8.7, 2.3, 3.6, 3.8, [("Gradient Boosting 0,756", {"size": 15, "bold": True, "color": TEAL}),
   ("LogReg 0,41 / SVM 0,39 / RF 0,71.", {"size": 13, "color": INK}), ("", {}),
   ("Les lineaires echouent : l'attaque a une signature NON LINEAIRE dans le CAN.", {"size": 13, "color": INK}),
   ("", {}), ("Biometrie seule = 0,014 (hasard).", {"size": 13, "color": MUTED})], space=6)

# ============================================================================ 13. CHEMIN B
s = slide(); title(s, "Chemin B : l'anomalie echoue (assume)",
                   sub="Apprendre le normal ne suffit pas ici")
pic(s, "p4b_anomaly.png", 0.8, 2.1, 7.0)
card(s, 8.2, 2.1, 4.4, 4.0)
tb(s, 8.5, 2.4, 3.8, 3.4, [("~ 0,02 = hasard", {"size": 18, "bold": True, "color": RED}),
   ("Les 4 detecteurs (Isolation Forest, One-Class SVM, PCA, gaussienne) sont au niveau du "
    "hasard.", {"size": 13, "color": INK}), ("", {}),
   ("La variabilite inter-conducteur noie l'attaque : rare n'est pas aberrant.", {"size": 13, "color": INK})], space=6)

# ============================================================================ 14. CHEMIN C
s = slide(); title(s, "Chemin C : le deep ne gagne pas",
                   sub="Multi-graine : MLP vs GRU, l'ecart est du bruit")
pic(s, "p4c_deep.png", 0.8, 2.2, 6.2)
card(s, 7.4, 2.2, 5.2, 3.7)
tb(s, 7.7, 2.5, 4.6, 3.1, [("MLP 0,543 +/- 0,016", {"size": 15, "bold": True, "color": INK}),
   ("GRU 0,571 +/- 0,024 -> ecart NON significatif.", {"size": 14, "color": INK}), ("", {}),
   ("Aucun ne bat les arbres (0,756). En tabulaire avec peu de positifs, les arbres dominent.", {"size": 14, "color": INK}),
   ("", {}), ("Le deep n'est pas un label de qualite.", {"size": 14, "bold": True, "color": TEAL})], space=6)

# ============================================================================ 15. EVALUATION ROC/PR
s = slide(); title(s, "Evaluation : ROC vs PR (la demande du sujet)",
                   sub="Une ROC flatteuse, une PR honnete")
pic(s, "v1_roc_vs_pr.png", 0.8, 2.0, 7.2)
card(s, 8.4, 2.0, 4.2, 4.1)
tb(s, 8.7, 2.3, 3.6, 3.5, [("AUC-ROC 0,977", {"size": 18, "bold": True, "color": MUTED}),
   ("parait excellente : 98,5 % de negatifs ecrasent les faux positifs.", {"size": 13, "color": INK}),
   ("", {}), ("PR-AUC 0,735", {"size": 18, "bold": True, "color": RED}),
   ("la metrique HONNETE pour une attaque rare.", {"size": 13, "color": INK})], space=6)

# ============================================================================ 16. TUNING + LATENCE
s = slide(); title(s, "Optimisation et metriques IDS",
                   sub="Hyperparametres (sujet) + detection par episode")
stat(s, 0.9, 2.3, 3.4, "0,757->0,798", "tuning (+0,040)", TEAL)
stat(s, 4.9, 2.3, 3.4, "86 %", "episodes detectes", NAVY)
stat(s, 8.9, 2.3, 3.4, "~4 s", "latence mediane", AMBER)
pic(s, "v1_latency_episode.png", 1.8, 3.7, 9.6)
tb(s, 0.8, 6.7, 11.8, 0.5, "Le tuning aide a la marge (sous l'ecart-type) ; au niveau EPISODE, la "
   "detection est bien meilleure que par fenetre.", size=13, color=MUTED)

# ============================================================================ 17. GRADIENT AWARENESS
s = slide(); title(s, "Une generalisation BIMODALE : le gradient d'awareness",
                   sub="L'IDS detecte en partie la REACTION du conducteur")
pic(s, "p5b_group_analysis.png", 0.8, 2.0, 7.2)
card(s, 8.4, 2.0, 4.2, 4.3)
tb(s, 8.7, 2.3, 3.6, 3.8, [("G1 0,74 / G2 0,92 / G3 0,96", {"size": 15, "bold": True, "color": INK}),
   ("La detectabilite suit l'AWARENESS.", {"size": 13, "color": INK}), ("", {}),
   ("G3 prevenu se gare -> le regime chute -> facile a voir. G1 non averti ne reagit pas -> "
    "dur (0,46).", {"size": 13, "color": INK}), ("", {}),
   ("-> on detecte la reaction autant que l'injection.", {"size": 13, "bold": True, "color": RED})], space=5)

# ============================================================================ 18. DECOUVERTE
s = slide(); title(s, "La decouverte : isoler la signature d'injection",
                   sub="Le bus CAN0 se TAIT ~4 s apres l'attaque (tous groupes)")
pic(s, "v2_injection_signature.png", 0.8, 2.0, 7.4)
card(s, 8.6, 2.0, 4.0, 4.3, fill=NAVY)
tb(s, 8.9, 2.3, 3.4, 3.8, [("67 % -> 6,7 %", {"size": 22, "bold": True, "color": RED, "font": SERIF}),
   ("La couverture du bus CAN0 s'effondre pendant l'attaque, marche d'escalier a +4 s.", {"size": 13, "color": WHITE}),
   ("", {}), ("Effet DIRECT de l'injection (l'ELD sature le bus), present meme sur le Groupe 1 non "
   "averti.", {"size": 13, "color": ICE}), ("", {}),
   ("-> une signature de l'intrusion, distincte de la reaction.", {"size": 13, "bold": True, "color": WHITE})], space=5)

# ============================================================================ 19. CARACTERISATION
s = slide(); title(s, "Deux empreintes simultanees a l'onset",
                   sub="On separe, signal par signal, l'injection de la reaction")
pic(s, "v2_attack_fingerprint.png", 0.8, 2.0, 8.0)
card(s, 9.0, 2.0, 3.6, 4.3)
tb(s, 9.25, 2.3, 3.1, 3.8, [("Disponibilite", {"size": 15, "bold": True, "color": RED}),
   ("seul le bus CAN0 se tait = l'INJECTION.", {"size": 13, "color": INK}), ("", {}),
   ("Valeur", {"size": 15, "bold": True, "color": NAVY}),
   ("la charge moteur (EGR, carburant, couple) s'effondre = la REACTION.", {"size": 13, "color": INK})], space=6)

# ============================================================================ 20. EVASION
s = slide(); title(s, "Robustesse (1/2) : l'attaquant adaptatif",
                   sub="Un adversaire qui connait le modele l'evade facilement")
pic(s, "v2_evasion.png", 0.8, 2.0, 7.0)
card(s, 8.2, 2.0, 4.4, 4.3, fill=NAVY)
tb(s, 8.5, 2.3, 3.8, 3.8, [("0,74 -> 0,22", {"size": 24, "bold": True, "color": RED, "font": SERIF}),
   ("Neutraliser UN seul signal (le CAN0) effondre la PR-AUC.", {"size": 13, "color": WHITE}),
   ("Deux signaux -> 0,07.", {"size": 13, "color": ICE}), ("", {}),
   ("Le detecteur est FRAGILE : son importance est ultra-concentree.", {"size": 13, "color": WHITE})], space=6)

# ============================================================================ 21. TAXONOMIE
s = slide(); title(s, "Robustesse (2/2) : detecteur mono-attaque",
                   sub="Le champion ne voit que l'attaque qu'on lui a montree")
pic(s, "v2_taxonomy.png", 0.8, 2.1, 7.0)
card(s, 8.2, 2.1, 4.4, 4.0)
tb(s, 8.5, 2.4, 3.8, 3.4, ["Attaque reelle : 65 %.", "",
   ("Fuzzing 0,7 % / Masquerade 0,9 % / Replay 1,1 % = ANGLES MORTS.", {"size": 14, "bold": True, "color": RED}),
   ("", {}), ("Seul le DoS/silence (15 %) ressemble a la signature apprise.", {"size": 13, "color": INK}),
   ("", {}), ("Perimetre valide = l'attaque ELD du dataset.", {"size": 13, "color": MUTED})], size=14, color=INK, space=6)

# ============================================================================ 22. DEMO
s = slide(dark=True)
title(s, "Place a la demonstration", dark=True,
      sub="Un IDS ne se juge pas sur une slide - il se manipule")
tb(s, 0.7, 1.6, 12, 0.7, "Jusqu'ici : la methode et les chiffres. Maintenant, le BANC D'ESSAI - "
   "vous pilotez le vrai modele, en direct.", size=16, color=ICE)
postes = [("1", "Voir l'IDS detecter", "une attaque reelle se deroule ; le bus CAN0 se tait sous vos yeux", TEAL),
          ("2", "Jouer l'attaquant", "neutralisez des signaux : l'IDS devient aveugle", RED),
          ("3", "Inventer une attaque", "DoS, fuzzing, masquerade : trouvez ses angles morts", RED),
          ("4", "Le piege methodo", "choisissez split + features : score honnete, ou mirage ?", AMBER),
          ("5", "Le deploiement", "seuil + rarete reelle : combien de fausses alertes ?", TEAL)]
y = 2.55
for num, head, txt, col in postes:
    circle(s, 0.9, y, 0.5, col, num, gsize=18)
    tb(s, 1.6, y - 0.02, 4.7, 0.5, head, size=17, bold=True, color=WHITE)
    tb(s, 6.1, y + 0.03, 6.5, 0.5, txt, size=13.5, color=ICE)
    y += 0.78
tb(s, 0.7, 6.75, 12, 0.5, "A vous : proposez vos scenarios, on les teste en direct.",
   size=15, color=WHITE, font=SERIF)

# ============================================================================ 23. AUTO-CRITIQUE
s = slide(); title(s, "Auto-critique et limites",
                   sub="Un projet honnete nomme ce qu'il n'a pas prouve")
card(s, 0.8, 2.0, 5.7, 4.3)
tb(s, 1.1, 2.3, 5.1, 0.5, "Traite au fil du projet", size=16, bold=True, color=TEAL)
tb(s, 1.1, 2.9, 5.1, 3.2, ["Confondeurs lieu/temps neutralises", "Fuite par conducteur chiffree",
   "ROC + tuning (demandes du sujet)", "Signature d'injection isolee (CAN0)",
   "Attaquant adaptatif, taxonomie", "Tests, versions, multi-graine"], size=13.5, color=INK, space=9)
card(s, 6.8, 2.0, 5.7, 4.3)
tb(s, 7.1, 2.3, 5.1, 0.5, "Limites residuelles", size=16, bold=True, color=RED)
tb(s, 7.1, 2.9, 5.1, 3.2, ["Cible conflee : injection + reaction",
   "Detecteur evadable (1-2 signaux)", "Mono-attaque (fuzzing/replay rates)",
   "Groupe 1 non averti : angle mort", "Generalisation hors-dataset non faite"], size=13.5, color=INK, space=9)

# ============================================================================ 23b. POSITIONNEMENT
s = slide(); title(s, "Positionnement & perspectives",
                   sub="Ce que dit la litterature, et ce qu'on ferait avec plus de temps")
card(s, 0.8, 2.0, 5.7, 4.3)
tb(s, 1.1, 2.3, 5.1, 0.5, "Positionnement", size=16, bold=True, color=TEAL)
tb(s, 1.1, 2.9, 5.1, 3.2, [
    "Papier source ORNL (2025) : valide notre gradient d'awareness, mais ne construit pas d'IDS.",
    "CIDS / CANet / ROAD : IDS CAN de reference, mais sur la TRAME brute (nous : agrege 1 s).",
    "Aucune PR-AUC de detection publiee ici -> cadrage IDS inedit, pas de SOTA direct."], size=13, color=INK, space=8)
card(s, 6.8, 2.0, 5.7, 4.3)
tb(s, 7.1, 2.3, 5.1, 0.5, "Perspectives (hors-scope)", size=16, bold=True, color=RED)
tb(s, 7.1, 2.9, 5.1, 3.2, [
    "Generalisation : rejouer sur ROAD / Car-Hacking.",
    "Defense en profondeur : 2e couche ORTHOGONALE (coherence physique, empreinte d'horloge).",
    "Donnees haute-resolution : isoler l'injection sous-seconde.",
    "Modele de menace elargi : entrainer sur une taxonomie d'attaques."], size=13, color=INK, space=8)

# ============================================================================ 24. CONCLUSION
s = slide(dark=True); title(s, "Conclusion", dark=True,
                            sub="D'un score apparent a une comprehension d'ingenieur securite")
points = [("1", "La vraie difficulte etait methodologique : les confondeurs et la fuite par conducteur."),
          ("2", "On peut isoler la signature de l'injection (silence CAN0), distincte de la reaction."),
          ("3", "L'honnetete - negatifs assumes, attaquant adaptatif, base rate - fait la credibilite.")]
y = 2.3
for num, txt in points:
    circle(s, 0.9, y, 0.6, RED, num, gsize=22)
    tb(s, 1.8, y, 10.6, 1.0, txt, size=18, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.25
tb(s, 0.9, 6.4, 11.5, 0.7, "Un IDS intelligent : un systeme qui detecte, et reste lucide sur ses "
   "incertitudes.", size=15, color=ICE, font=SERIF)

# ============================================================================ 25. MERCI
s = slide(dark=True)
tb(s, 0.9, 3.0, 11.5, 1.2, "Merci de votre attention", size=40, bold=True, font=SERIF, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.9, 4.3, 11.5, 0.6, "Questions et discussion", size=18, color=ICE, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"[OK] presentation -> {OUT}  ({len(prs.slides._sldIdLst)} diapositives)")
